from fastapi import FastAPI, UploadFile, File, HTTPException, status, BackgroundTasks
from fastapi.responses import JSONResponse
from pathlib import Path
import os
import aiofiles
import pandas as pd
import logging
import asyncio
from typing import Optional
import uuid
from datetime import datetime
import json

# Docx and Pptx processing libraries
from docx import Document
from pptx import Presentation

# Set up logging
logging.basicConfig(level=logging.INFO, format='%(asctime)s - %(levelname)s - %(message)s')
logger = logging.getLogger(__name__)

# FastAPI app instance
app = FastAPI(
    title="Large File Upload Service",
    description="Service for uploading and processing large files (Excel, Word, PPT) to JSON",
    version="1.0.0"
)

# Configuration
CHUNK_SIZE = 1024 * 1024  # 1MB chunks
ALLOWED_EXTENSIONS = {".xlsx", ".xls", ".pptx", ".ppt", ".docx", ".doc"}
EXCEL_EXTENSIONS = {".xlsx", ".xls"}
WORD_EXTENSIONS = {".docx", ".doc"}
PPT_EXTENSIONS = {".pptx", ".ppt"}

# Create directories
UPLOAD_DIR = Path("uploads")
PROCESSED_DATA_DIR = Path("processed_data")
UPLOAD_DIR.mkdir(exist_ok=True)
PROCESSED_DATA_DIR.mkdir(exist_ok=True)

# In-memory task tracking (use Redis/database in production)
task_status = {}

# --- Helper functions ---

async def save_upload_file_chunked(upload_file: UploadFile, destination_path: Path) -> int:
    """Saves an UploadFile to disk in chunks and returns the total bytes written."""
    total_bytes = 0
    try:
        async with aiofiles.open(destination_path, "wb") as out_file:
            while True:
                chunk = await upload_file.read(CHUNK_SIZE)
                if not chunk:
                    break
                total_bytes += len(chunk)
                await out_file.write(chunk)
        logger.info(f"File saved successfully: {destination_path} ({total_bytes} bytes)")
        return total_bytes
    except Exception as e:
        logger.error(f"Error saving file in chunks: {e}")
        if destination_path.exists():
            destination_path.unlink()
        raise HTTPException(status_code=500, detail="Failed to save file")
    finally:
        await upload_file.close()

def update_task_status(task_id: str, status: str, details: Optional[str] = None):
    """Update task status in memory (use database in production)."""
    task_status[task_id] = {
        "status": status,
        "details": details,
        "timestamp": datetime.utcnow().isoformat()
    }

def get_word_file_content(file_path: Path) -> list:
    """Extracts text and tables from a Word file."""
    content = []
    try:
        doc = Document(file_path)
        for i, paragraph in enumerate(doc.paragraphs):
            if paragraph.text.strip():
                content.append({"type": "paragraph", "text": paragraph.text, "index": i})
        
        for i, table in enumerate(doc.tables):
            table_data = []
            for row_idx, row in enumerate(table.rows):
                row_data = [cell.text for cell in row.cells]
                table_data.append(row_data)
            content.append({"type": "table", "data": table_data, "index": i + len(doc.paragraphs)})
    except Exception as e:
        logger.error(f"Error processing Word file: {e}")
        raise
    return content

def get_ppt_file_content(file_path: Path) -> list:
    """Extracts text from a PowerPoint file."""
    content = []
    try:
        prs = Presentation(file_path)
        for i, slide in enumerate(prs.slides):
            slide_content = {"slide_number": i + 1, "elements": []}
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        if paragraph.text.strip():
                            slide_content["elements"].append({"type": "paragraph", "text": paragraph.text})
                if shape.has_table:
                    table_data = []
                    for row in shape.table.rows:
                        row_data = [cell.text for cell in row.cells]
                        table_data.append(row_data)
                    slide_content["elements"].append({"type": "table", "data": table_data})
            
            if slide_content["elements"]:
                content.append(slide_content)
    except Exception as e:
        logger.error(f"Error processing PPT file: {e}")
        raise
    return content

async def process_file_background(file_path: Path, original_filename: str, task_id: str):
    """Process a file in the background and save JSON output."""
    try:
        update_task_status(task_id, "processing", "Starting file processing")
        if not file_path.exists():
            raise FileNotFoundError(f"File not found: {file_path}")

        file_extension = Path(original_filename).suffix.lower()
        processed_data = None
        
        # Dispatch to the correct processing function based on file extension
        if file_extension in EXCEL_EXTENSIONS:
            update_task_status(task_id, "processing", "Reading Excel file")
            all_sheets = pd.read_excel(file_path, sheet_name=None)
            combined_data = []
            for sheet_name, sheet_df in all_sheets.items():
                if not sheet_df.empty:
                    sheet_df["source_sheet"] = sheet_name
                    combined_data.append(sheet_df)
            
            if combined_data:
                df = pd.concat(combined_data, ignore_index=True, sort=False)
                df["processed_at"] = datetime.utcnow().isoformat()
                processed_data = df.to_dict(orient="records")
            else:
                update_task_status(task_id, "failed", "No valid data found in Excel")
                return
        
        elif file_extension in WORD_EXTENSIONS:
            update_task_status(task_id, "processing", "Reading Word file")
            processed_data = get_word_file_content(file_path)
        
        elif file_extension in PPT_EXTENSIONS:
            update_task_status(task_id, "processing", "Reading PowerPoint file")
            processed_data = get_ppt_file_content(file_path)
            
        if processed_data is None:
            update_task_status(task_id, "failed", f"Failed to process content from {file_extension} file")
            return

        processed_filename = f"processed_{Path(original_filename).stem}_{task_id}.json"
        processed_filepath = PROCESSED_DATA_DIR / processed_filename
        
        # Write JSON output
        with open(processed_filepath, "w") as f:
            json.dump(processed_data, f, indent=2)

        update_task_status(task_id, "completed", f"Output saved to {processed_filename}")
    
    except Exception as e:
        logger.error(f"Processing task failed for {task_id}: {e}")
        update_task_status(task_id, "failed", str(e))
    finally:
        if file_path.exists():
            file_path.unlink()

# --- Endpoints ---

@app.post("/upload/to-json")
async def upload_file_to_json(file: UploadFile = File(...), background_tasks: BackgroundTasks = BackgroundTasks()):
    """
    Upload a large file (Excel, Word, PPT) and process it in the background, outputting JSON.
    The file is uploaded in chunks to save memory on the server.
    """
    if not file.filename:
        raise HTTPException(status_code=400, detail="No file provided")

    file_extension = Path(file.filename).suffix.lower()
    if file_extension not in ALLOWED_EXTENSIONS:
        raise HTTPException(status_code=400, detail=f"Only {', '.join(ALLOWED_EXTENSIONS)} files are allowed")

    task_id = str(uuid.uuid4())
    unique_filename = f"{Path(file.filename).stem}_{task_id}{file_extension}"
    destination_path = UPLOAD_DIR / unique_filename

    update_task_status(task_id, "uploading", "Saving file to disk in chunks")
    file_size = await save_upload_file_chunked(file, destination_path)

    # Add the processing task to the background
    background_tasks.add_task(process_file_background, destination_path, file.filename, task_id)
    update_task_status(task_id, "queued", "File uploaded successfully, processing queued for JSON output")

    return {
        "task_id": task_id,
        "filename": file.filename,
        "file_size": file_size,
        "output_format": "json",
        "message": "File uploaded successfully and queued for JSON processing",
        "status_url": f"/status/{task_id}"
    }

@app.get("/status/{task_id}")
async def get_task_status(task_id: str):
    """Retrieve the status of a specific background processing task."""
    if task_id not in task_status:
        raise HTTPException(status_code=404, detail="Task not found")
    return {"task_id": task_id, **task_status[task_id]}

@app.get("/health")
async def health_check():
    """Endpoint to check the service's health status."""
    return {"status": "healthy", "timestamp": datetime.utcnow().isoformat()}

# Error handler
@app.exception_handler(HTTPException)
async def http_exception_handler(request, exc):
    return JSONResponse(status_code=exc.status_code, content={"detail": exc.detail})

if __name__ == "__main__":
    import uvicorn
    uvicorn.run(app, host="0.0.0.0", port=8000)
